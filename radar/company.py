from __future__ import annotations

import io
from pathlib import Path

import requests
from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from radar.db import active_company, rows

MAX_DOCUMENT_CHARS = 80_000
MAX_CONTEXT_CHARS = 24_000
USER_AGENT = "Innovation-Radar/0.2 (student research prototype)"


class DocumentError(ValueError):
    pass


def _pdf_text(data: bytes) -> str:
    return "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(data)).pages)


def _pptx_text(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    return "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))


def _docx_text(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def extract_document(name: str, data: bytes, content_type: str = "") -> str:
    suffix = Path(name).suffix.lower()
    try:
        if suffix == ".pdf" or "pdf" in content_type:
            text = _pdf_text(data)
        elif suffix == ".pptx" or "presentation" in content_type:
            text = _pptx_text(data)
        elif suffix == ".docx" or "wordprocessing" in content_type:
            text = _docx_text(data)
        elif suffix in {".txt", ".md", ".csv", ".json"} or content_type.startswith("text/"):
            text = data.decode("utf-8", errors="replace")
        elif "html" in content_type:
            text = BeautifulSoup(data, "html.parser").get_text(" ", strip=True)
        else:
            raise DocumentError(f"Unsupported document type: {suffix or content_type or 'unknown'}")
    except Exception as error:
        if isinstance(error, DocumentError):
            raise
        raise DocumentError(f"Could not extract {name}: {error}") from error
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not text:
        raise DocumentError(f"No readable text found in {name}. Scanned PDFs need OCR, which is not implemented yet.")
    return text[:MAX_DOCUMENT_CHARS]


def fetch_document_url(url: str) -> tuple[str, str]:
    response = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=30)
    response.raise_for_status()
    name = Path(response.url.split("?", 1)[0]).name or response.url
    return name, extract_document(name, response.content, response.headers.get("Content-Type", ""))


def company_context(max_chars: int = MAX_CONTEXT_CHARS) -> str:
    company = active_company()
    documents = rows("SELECT name,source_type,source_url,extracted_text FROM company_documents ORDER BY added_at DESC")
    parts = [
        f"Active company: {company.get('name', 'Not configured')}",
        f"Priority geography: {company.get('geography', '')}",
        f"Website: {company.get('website_url', '')}",
        f"Company-specific research instruction: {company.get('strategic_prompt', '')}",
        "Partner rule: consider direct, partner-led, ecosystem, reseller, integrator, and capability-gap opportunities. External delivery is not automatically negative; explain the company's role and dependency.",
    ]
    for document in documents:
        parts.append(f"\nREFERENCE DOCUMENT: {document['name']} ({document['source_type']}, {document['source_url']})\n{document['extracted_text']}")
    return "\n".join(parts)[:max_chars]
