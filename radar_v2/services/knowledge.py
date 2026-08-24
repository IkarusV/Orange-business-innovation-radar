from __future__ import annotations

import io
import json
from pathlib import Path

from docx import Document
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation

from radar_v2.constants import DOCUMENTS
from radar_v2.services import extension_store


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    data = path.read_bytes()
    if suffix == ".pdf":
        text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(data)).pages)
    elif suffix == ".pptx":
        presentation = Presentation(io.BytesIO(data))
        text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    elif suffix == ".docx":
        document = Document(io.BytesIO(data))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    elif suffix in {".txt", ".md", ".csv", ".json", ".html", ".htm"}:
        text = data.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")
    cleaned = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not cleaned:
        raise ValueError("No readable text was found in this document")
    return cleaned[:100_000]


def _client(base_url: str, api_key: str) -> OpenAI:
    if not api_key:
        raise ValueError("AI provider key is required")
    return OpenAI(api_key=api_key, base_url=base_url.rstrip("/"))


def _json_result(client: OpenAI, model: str, mode: str, system: str, content: str) -> dict:
    if mode == "chat":
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "system", "content": system}, {"role": "user", "content": content}],
            response_format={"type": "json_object"},
        )
        raw = response.choices[0].message.content or "{}"
    else:
        response = client.responses.create(model=model, instructions=system, input=content)
        raw = response.output_text
    raw = raw.strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
    return json.loads(raw)


def process_document(document: dict, instruction: str, base_url: str, api_key: str, model: str, mode: str) -> Path:
    raw_path = Path(document["path"])
    text = extract_text(raw_path)
    system = """You are a company knowledge analyst. Summarize only the supplied document; never use facts from another document or outside knowledge. Return JSON with executive_summary, key_facts, financial_signals, strategic_priorities, company_vocabulary, capabilities, risks_and_unknowns, and useful_radar_guidance. Mark uncertainty and keep numbers with their units and periods."""
    user_focus = instruction.strip() or "Provide a balanced company-relevance summary."
    result = _json_result(
        _client(base_url, api_key), model, mode, system,
        f"DOCUMENT: {document['name']}\nUSER FOCUS: {user_focus}\n\nDOCUMENT TEXT:\n{text}",
    )
    company = extension_store.active_company()
    folder = DOCUMENTS / extension_store.safe_name(company["name"]) / "processed"
    folder.mkdir(parents=True, exist_ok=True)
    target = folder / f"{raw_path.stem}.summary.json"
    target.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    extension_store.update_document_processing(document["id"], "Processed", str(target), instruction)
    return target


def create_combined_report(documents: list[dict], instruction: str, base_url: str, api_key: str, model: str, mode: str) -> Path:
    if len(documents) < 2:
        raise ValueError("Select at least two documents")
    sections = []
    for document in documents:
        source_path = Path(document.get("processed_path") or document["path"])
        text = source_path.read_text(encoding="utf-8", errors="replace") if source_path.suffix.lower() in {".json", ".txt", ".md"} else extract_text(source_path)
        sections.append(f"SOURCE: {document['name']}\n{text[:35_000]}")
    system = """You are a senior company research analyst. Synthesize the selected company documents into one decision-ready report. Preserve source boundaries, identify repeated themes and contradictions, separate facts from interpretations, and return JSON with report_summary, financial_profile, strategic_priorities, capabilities, company_vocabulary, repeated_themes, contradictions, opportunity_preferences, partnership_preferences, risks_and_unknowns, and radar_guidance."""
    result = _json_result(
        _client(base_url, api_key), model, mode, system,
        f"USER FOCUS: {instruction.strip() or 'Create a balanced company profile.'}\n\n" + "\n\n".join(sections),
    )
    company = extension_store.active_company()
    folder = DOCUMENTS / extension_store.safe_name(company["name"]) / "processed"
    folder.mkdir(parents=True, exist_ok=True)
    index = len(list(folder.glob("company_report_*.json"))) + 1
    target = folder / f"company_report_{index}.json"
    target.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    extension_store.save_company_report(f"{company['name']} knowledge report {index}", target, len(documents))
    return target
